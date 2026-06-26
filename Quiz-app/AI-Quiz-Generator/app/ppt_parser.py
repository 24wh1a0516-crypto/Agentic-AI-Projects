from pathlib import Path
from typing import List

from pptx import Presentation


def extract_text_from_ppt(file_path: str) -> List[str]:
    path = Path(file_path)
    if not path.exists():
        raise FileNotFoundError(f"File not found: {file_path}")

    presentation = Presentation(path)
    slides_text = []

    for slide in presentation.slides:
        text_chunks = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                text_chunks.append(shape.text.strip())
        slides_text.append("\n".join(text_chunks))

    return slides_text
